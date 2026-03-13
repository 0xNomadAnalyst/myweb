"""
Generate an Opus-branded PowerPoint theme (.pptx) matching the opus landing page
visual identity: dark surface hierarchy, orange accent, Inter / Consolas typography.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy
import os

# ── Opus color palette ──────────────────────────────────────────────────────
BG_L0 = RGBColor(0x0B, 0x12, 0x20)      # --background  page base
BG_L1 = RGBColor(0x0F, 0x17, 0x2A)      # --surface     section bands
BG_L2 = RGBColor(0x11, 0x1C, 0x34)      # --secondary   containers / tags
BG_L3 = RGBColor(0x17, 0x20, 0x33)      # --card        cards / panels
BG_L4 = RGBColor(0x1E, 0x29, 0x3B)      # --accent      hover / interactive
PRIMARY = RGBColor(0xFF, 0x6B, 0x00)     # --primary     orange accent
CTA = RGBColor(0xF8, 0xA9, 0x4A)        # --cta         warm gold accent
FG = RGBColor(0xE5, 0xE7, 0xEB)         # --foreground  primary text
MUTED = RGBColor(0x9C, 0xA3, 0xAF)      # --muted-foreground
BORDER = RGBColor(0x1F, 0x27, 0x37)      # approximate oklch(1 0 0 / 6%) on dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)  # 16:9 at 96 dpi
SLIDE_H = Inches(7.5)

FONT_HEADING = "Inter"
FONT_BODY = "Inter"
FONT_MONO = "Consolas"


def hex_to_srgb(r, g, b):
    return f"{r:02X}{g:02X}{b:02X}"


def set_slide_bg(slide, color):
    """Set solid background fill on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_HEADING,
                font_size=Pt(18), font_color=FG, bold=False, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    tf.auto_size = None

    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }
        bodyPr.set("anchor", anchor_map.get(anchor, "t"))

    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = Pt(1)
    else:
        ln.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=BORDER, width=Pt(0.75)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_bullet_list(slide, left, top, width, height, items, font_size=Pt(14)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = font_size
        p.font.color.rgb = MUTED
        p.space_after = Pt(8)
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        # Remove existing bullet elements
        for child in list(pPr):
            if child.tag in (qn("a:buFont"), qn("a:buChar"), qn("a:buNone")):
                pPr.remove(child)
        pPr.append(buFont)
        pPr.append(buChar)
    return txBox


def add_tag_pill(slide, left, top, text):
    w, h = Inches(1.8), Inches(0.35)
    shape = add_rect(slide, left, top, w, h, BG_L2, BORDER)
    shape.text = ""
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_MONO
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", "ctr")
    return shape


def build_theme():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Title
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.6),
                "Operational Intelligence for\nOnchain Financial Systems",
                font_size=Pt(40), font_color=FG, bold=True,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.9),
                "Analytical infrastructure, real-time monitoring, and risk modelling\n"
                "for teams operating complex digital-asset environments.",
                font_size=Pt(16), font_color=MUTED,
                alignment=PP_ALIGN.CENTER)

    add_line(slide, Inches(5.5), Inches(4.8), Inches(7.8), Inches(4.8),
             color=MUTED, width=Pt(0.5))

    add_textbox(slide, Inches(3.5), Inches(5.1), Inches(6.3), Inches(0.4),
                "Roderick McKinley, CFA, FRM",
                font_size=Pt(13), font_color=MUTED,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(0.4),
                "Independent Financial Systems Analyst",
                font_size=Pt(11), font_color=RGBColor(0x6B, 0x72, 0x80),
                font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 2: The Operational Reality
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10.9), Inches(0.7),
                "Onchain Financial Systems Operate in Continuous Time",
                font_size=Pt(30), font_color=FG, bold=True)

    paras = [
        "Digital-asset markets evolve continuously: liquidity conditions, leverage, and execution outcomes shift in real time.",
        "While blockchain networks make these dynamics publicly observable, decision-grade visibility typically requires substantial analytical and infrastructure development.",
        "Teams operating these systems need continuous insight into liquidity, execution conditions, protocol exposures, and emerging stress signals.",
    ]
    y = Inches(2.0)
    for para in paras:
        add_textbox(slide, Inches(1.2), y, Inches(10.5), Inches(0.8),
                    para, font_size=Pt(15), font_color=MUTED)
        y += Inches(1.2)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 3: Operational Visibility Domains
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.7), Inches(10.9), Inches(0.6),
                "Where Visibility Is Required",
                font_size=Pt(30), font_color=FG, bold=True)

    domains = [
        ("DEX Liquidity & Pricing (CLMMs)",
         "Liquidity structure across price ranges and price impact around the current price."),
        ("Lending & Collateral Risk (e.g., Kamino)",
         "Utilisation, liquidation conditions, and collateral concentration dynamics."),
        ("Yield / Structured Exposure (e.g., Exponent)",
         "Vault state, rate dynamics, and position behaviour."),
        ("Cross-Protocol Stress Signals",
         "Dependencies and stress propagation across venues and protocols."),
        ("Partner Oversight (Market Makers / LPs)",
         "Independent evaluation of depth provision and behaviour during volatility."),
    ]

    cols = 3
    card_w, card_h = Inches(3.5), Inches(1.8)
    gap = Inches(0.3)
    start_x = Inches(1.2)
    start_y = Inches(1.7)

    for i, (title, desc) in enumerate(domains):
        row, col = divmod(i, cols)
        if row == 1 and len(domains) - cols < cols:
            remaining = len(domains) - cols
            offset = (cols - remaining) * (card_w + gap) / 2
            x = start_x + col * (card_w + gap) + int(offset)
        else:
            x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        add_rect(slide, x, y, card_w, card_h, BG_L3, BORDER)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.5), Inches(0.35),
                    title, font_size=Pt(13), font_color=FG, bold=True)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.6), card_w - Inches(0.5), Inches(1.0),
                    desc, font_size=Pt(11), font_color=MUTED)

    add_textbox(slide, Inches(1.2), Inches(5.9), Inches(10.5), Inches(0.4),
                "Turnkey coverage today: Solana DEX + Kamino + Exponent. "
                "Extendable to additional Solana/EVM protocols as required.",
                font_size=Pt(10), font_color=RGBColor(0x6B, 0x72, 0x80),
                font_name=FONT_MONO)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 4: Flagship — Realtime DeFi Monitoring Station
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(5), Inches(0.3),
                "FLAGSHIP CAPABILITY", font_name=FONT_MONO,
                font_size=Pt(9), font_color=MUTED)
    add_textbox(slide, Inches(1.2), Inches(0.9), Inches(10.9), Inches(0.6),
                "Realtime DeFi Monitoring Station",
                font_size=Pt(28), font_color=FG, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.5),
                "Continuous operational visibility across pricing, exposure, "
                "and counterparties \u2014 in one place.",
                font_size=Pt(14), font_color=MUTED)

    bullets = [
        "Unified operational picture across DEX pricing/liquidity, lending exposure, and yield positions.",
        "Decision-grade monitoring built around actions ops teams can actually take.",
        "Cross-protocol risk surfaces that aggregate signal across venues.",
        "Historical event store for ad hoc investigation and post-stress diagnosis.",
        "Extensible architecture: turnkey Solana modules today; rapid build-out as needed.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(7), Inches(3.5), bullets)

    # Right-side inset panels
    panel_x = Inches(8.8)
    add_rect(slide, panel_x, Inches(2.5), Inches(3.8), Inches(1.4), BG_L4, BORDER)
    add_textbox(slide, panel_x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(0.25),
                "TURNKEY COVERAGE (TODAY)", font_name=FONT_MONO,
                font_size=Pt(8), font_color=MUTED)
    add_textbox(slide, panel_x + Inches(0.3), Inches(3.0), Inches(3.2), Inches(0.6),
                "Solana DEX monitoring\n+ Kamino + Exponent",
                font_size=Pt(13), font_color=FG)

    add_rect(slide, panel_x, Inches(4.2), Inches(3.8), Inches(1.2), BG_L4, BORDER)
    add_textbox(slide, panel_x + Inches(0.3), Inches(4.3), Inches(3.2), Inches(0.25),
                "EXTENDABLE", font_name=FONT_MONO,
                font_size=Pt(8), font_color=MUTED)
    add_textbox(slide, panel_x + Inches(0.3), Inches(4.7), Inches(3.2), Inches(0.5),
                "Additional Solana/EVM protocols",
                font_size=Pt(13), font_color=FG)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 5: Differentiator — Independent MM Monitoring
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(5), Inches(0.3),
                "DIFFERENTIATOR", font_name=FONT_MONO,
                font_size=Pt(9), font_color=MUTED)
    add_textbox(slide, Inches(1.2), Inches(0.9), Inches(10.9), Inches(0.6),
                "Independent Market Maker Monitoring",
                font_size=Pt(28), font_color=FG, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.5),
                "Objective oversight of third-party liquidity support \u2014 "
                "beyond reports, promises, and UI-level impressions.",
                font_size=Pt(14), font_color=MUTED)

    bullets = [
        "Verify performance against obligations: quantify depth provision and price-impact conditions.",
        "Stress-period behaviour matters most: detect liquidity withdrawal or instability when markets move.",
        "Comparable metrics across venues: consistent measures rather than ad hoc screenshots.",
        "Incentive alignment & contract improvement: KPI definitions, better terms, grounded in observed behaviour.",
        "Replace weak partnerships: identify inadequate support early, enable transitions.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(3.0), bullets)

    enables = ["renegotiation", "KPI redesign", "incentive alignment",
               "provider replacement", "in-house build decisions"]
    add_textbox(slide, Inches(1.2), Inches(5.5), Inches(3), Inches(0.3),
                "ENABLES", font_name=FONT_MONO, font_size=Pt(8), font_color=MUTED)
    pill_x = Inches(1.2)
    for tag in enables:
        add_tag_pill(slide, pill_x, Inches(5.9), tag)
        pill_x += Inches(2.0)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 6: Cross-Protocol Monitoring Station
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10.9), Inches(0.6),
                "Cross-Protocol Monitoring Station",
                font_size=Pt(28), font_color=FG, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(0.5),
                "As stablecoins and yield-bearing RWAs span multiple DeFi services, "
                "risk becomes system-level \u2014 visibility must too.",
                font_size=Pt(14), font_color=MUTED)

    bullets = [
        "Unified cross-protocol view: consolidate key signals into one operational picture.",
        "Normalized metrics across protocols: comparable definitions, consistent measurement.",
        "System-level risk surfaces: track dependencies between liquidity, leverage, and execution.",
        "Action-oriented dashboards: indicators aligned with interventions teams can actually take.",
        "Extendable coverage: turnkey Solana modules today; expand to additional protocols as needed.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(3.5), bullets)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Incident Replay & Root-Cause Analysis
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(5), Inches(0.3),
                "PREMIUM WORKFLOW", font_name=FONT_MONO,
                font_size=Pt(9), font_color=MUTED)
    add_textbox(slide, Inches(1.2), Inches(0.9), Inches(10.9), Inches(0.6),
                "Incident Replay & Root-Cause Analysis",
                font_size=Pt(28), font_color=FG, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.7),
                "A historical event store turns stress events into measurable lessons \u2014 "
                "reducing time-to-understanding and accelerating policy refinement.",
                font_size=Pt(14), font_color=MUTED)

    bullets = [
        "Historical event store across covered protocols for fast investigation.",
        "Rapid ad hoc diagnostics: shorten the path from \u201csomething happened\u201d to \u201cwe understand why.\u201d",
        "Root-cause analysis across venues: identify mechanisms, not just correlations.",
        "Accelerate the learning loop: refine playbooks, policies, and thresholds.",
        "Supports governance and reporting: clear narratives grounded in measurable evidence.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(10.5), Inches(3.0), bullets)

    add_textbox(slide, Inches(1.2), Inches(5.5), Inches(3), Inches(0.3),
                "OUTCOME", font_name=FONT_MONO, font_size=Pt(8), font_color=MUTED)
    outcomes = ["faster recovery", "fewer repeats", "better policies"]
    pill_x = Inches(1.2)
    for tag in outcomes:
        add_tag_pill(slide, pill_x, Inches(5.9), tag)
        pill_x += Inches(2.0)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 8: Simulation & Risk Policy Modelling
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(5), Inches(0.3),
                "FORWARD-LOOKING CAPABILITY", font_name=FONT_MONO,
                font_size=Pt(9), font_color=MUTED)
    add_textbox(slide, Inches(1.2), Inches(0.9), Inches(10.9), Inches(0.6),
                "Simulation & Risk Policy Modelling",
                font_size=Pt(28), font_color=FG, bold=True)
    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(0.5),
                "Evaluate policies and new design choices quantitatively "
                "before they become operational risk.",
                font_size=Pt(14), font_color=MUTED)

    bullets = [
        "Risk policy design & evaluation: model thresholds, provisions, and playbooks against history.",
        "Liquidity provisioning & capacity modelling: estimate sizing for buffers and internal operations.",
        "Feature / mechanism cost analysis: evaluate cost and risk impact before deployment.",
        "Stress testing and scenario analysis: explore outcomes under adverse conditions.",
        "Grounded in observed data: inputs from empirical monitoring, not theoretical assumptions.",
    ]
    add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(10.5), Inches(3.0), bullets)

    add_textbox(slide, Inches(1.2), Inches(5.5), Inches(3), Inches(0.3),
                "OUTPUTS", font_name=FONT_MONO, font_size=Pt(8), font_color=MUTED)
    outputs = ["policy recommendations", "parameter ranges",
               "sizing estimates", "stress scenarios"]
    pill_x = Inches(1.2)
    for tag in outputs:
        add_tag_pill(slide, pill_x, Inches(5.9), tag)
        pill_x += Inches(2.0)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 9: Close — Engagement Model
    # ────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_L0)

    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(5), Inches(0.3),
                "ENGAGEMENT MODEL", font_name=FONT_MONO,
                font_size=Pt(9), font_color=MUTED)
    add_textbox(slide, Inches(1.2), Inches(0.9), Inches(10.9), Inches(0.6),
                "Principal-Led",
                font_size=Pt(28), font_color=FG, bold=True)

    principles = [
        ("Principal-led delivery",
         "Direct work with the practitioner defining analysis, models, and systems \u2014 "
         "no delivery layers or handoffs."),
        ("Integrated advisory + build",
         "Empirical analysis and implementation stay aligned, preserving decision clarity end-to-end."),
        ("Designed for internal ownership",
         "Monitoring environments and analytical frameworks transferred so teams can "
         "operate and extend independently."),
    ]

    card_w = Inches(3.5)
    card_h = Inches(1.8)
    gap = Inches(0.3)
    start_x = Inches(1.2)
    y_cards = Inches(1.8)
    for i, (title, desc) in enumerate(principles):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, y_cards, card_w, card_h, BG_L3, BORDER)
        add_textbox(slide, x + Inches(0.25), y_cards + Inches(0.2),
                    card_w - Inches(0.5), Inches(0.35),
                    title, font_size=Pt(13), font_color=FG, bold=True)
        add_textbox(slide, x + Inches(0.25), y_cards + Inches(0.65),
                    card_w - Inches(0.5), Inches(1.0),
                    desc, font_size=Pt(11), font_color=MUTED)

    # Credentials
    cred_y = Inches(4.1)
    add_textbox(slide, Inches(1.2), cred_y, Inches(3), Inches(0.2),
                "BACKGROUND", font_name=FONT_MONO, font_size=Pt(8), font_color=MUTED)
    add_textbox(slide, Inches(1.2), cred_y + Inches(0.25), Inches(5), Inches(0.3),
                "Institutional finance + independent consulting",
                font_size=Pt(12), font_color=FG)
    add_textbox(slide, Inches(1.2), cred_y + Inches(0.7), Inches(3), Inches(0.2),
                "CREDENTIALS", font_name=FONT_MONO, font_size=Pt(8), font_color=MUTED)
    add_textbox(slide, Inches(1.2), cred_y + Inches(0.95), Inches(5), Inches(0.3),
                "CFA Charterholder  \u2022  FRM",
                font_size=Pt(12), font_color=FG)
    add_textbox(slide, Inches(1.2), cred_y + Inches(1.4), Inches(3), Inches(0.2),
                "CURRENT FOCUS", font_name=FONT_MONO, font_size=Pt(8), font_color=MUTED)
    add_textbox(slide, Inches(1.2), cred_y + Inches(1.65), Inches(6), Inches(0.3),
                "Operational intelligence for onchain financial systems",
                font_size=Pt(12), font_color=FG)

    # CTA / Contact
    add_textbox(slide, Inches(7.5), Inches(4.3), Inches(5), Inches(0.8),
                "If you operate a digital-asset system where visibility and "
                "execution quality are mission-critical, let\u2019s discuss scope and deployment.",
                font_size=Pt(13), font_color=MUTED, alignment=PP_ALIGN.RIGHT)

    contact_y = Inches(5.5)
    add_textbox(slide, Inches(7.5), contact_y, Inches(5), Inches(0.25),
                "Email: [your email]", font_name=FONT_MONO,
                font_size=Pt(10), font_color=RGBColor(0x6B, 0x72, 0x80),
                alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(7.5), contact_y + Inches(0.3), Inches(5), Inches(0.25),
                "Website: [your site]", font_name=FONT_MONO,
                font_size=Pt(10), font_color=RGBColor(0x6B, 0x72, 0x80),
                alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(7.5), contact_y + Inches(0.6), Inches(5), Inches(0.25),
                "LinkedIn: [your LinkedIn]", font_name=FONT_MONO,
                font_size=Pt(10), font_color=RGBColor(0x6B, 0x72, 0x80),
                alignment=PP_ALIGN.RIGHT)

    # ── Save ────────────────────────────────────────────────────────────────
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "opus-capabilities-deck.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    build_theme()
